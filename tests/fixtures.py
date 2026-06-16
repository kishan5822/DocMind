"""Programmatic test fixtures — generate small real files of each type in memory."""
from __future__ import annotations

import io
import json

import pandas as pd


def make_pdf(text_blocks) -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    y = 72
    for block in text_blocks:
        page.insert_text((72, y), block, fontsize=11)
        y += 18
    data = doc.tobytes()
    doc.close()
    return data


def make_docx(heading: str, paragraphs, table_rows=None) -> bytes:
    import docx

    d = docx.Document()
    d.add_heading(heading, level=1)
    for p in paragraphs:
        d.add_paragraph(p)
    if table_rows:
        t = d.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r, row in enumerate(table_rows):
            for c, val in enumerate(row):
                t.cell(r, c).text = str(val)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def make_pptx(slides) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[5]  # title only
    for title, body in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = title
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
        tb.text_frame.text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_xlsx(df: pd.DataFrame) -> bytes:
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False, sheet_name="Data")
    return buf.getvalue()


def make_csv(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8")


def make_json(obj) -> bytes:
    return json.dumps(obj).encode("utf-8")


def make_html(title: str, body: str) -> bytes:
    return f"<html><head><title>{title}</title></head><body><p>{body}</p></body></html>".encode("utf-8")


def make_txt(text: str) -> bytes:
    return text.encode("utf-8")
