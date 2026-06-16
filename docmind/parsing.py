"""Stage 2 — parsing.

Turns each accepted file into clean text organised into sections. A section
carries an optional heading (page number, slide title, sheet name, doc heading)
that the chunker uses to build the required metadata prefix.

PyMuPDF is the default PDF parser (fast, light). Docling is opt-in via
ENABLE_DOCLING for complex/scanned PDFs and is imported lazily so the base
install stays slim.
"""
from __future__ import annotations

import io
import json
from dataclasses import dataclass, field
from typing import List, Optional

from .config import config
from .logging_config import get_logger
from .validation import AcceptedFile

logger = get_logger(__name__)


class ParsingError(Exception):
    """Raised when a file cannot be parsed into usable text."""


@dataclass
class Section:
    """A logical span of text with an optional heading for metadata."""

    text: str
    heading: Optional[str] = None


@dataclass
class ParsedDocument:
    """Parsed result for one file."""

    filename: str
    sections: List[Section] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(s.text for s in self.sections if s.text.strip())

    @property
    def is_empty(self) -> bool:
        return not self.full_text.strip()


def parse_file(f: AcceptedFile) -> ParsedDocument:
    """Dispatch to the right parser by category. Raises ParsingError on failure."""
    dispatch = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
        "xlsx": _parse_xlsx,
        "csv": _parse_csv,
        "json": _parse_json,
        "text": _parse_text,
        "html": _parse_html,
        "image": _parse_image,
    }
    parser = dispatch.get(f.category)
    if parser is None:
        raise ParsingError(f"No parser for category '{f.category}'.")
    try:
        doc = parser(f)
    except ParsingError:
        raise
    except Exception as e:  # convert any parser crash into a clean error
        logger.exception("Parser crashed on '%s'", f.name)
        raise ParsingError(f"Failed to parse '{f.name}': {e}") from e

    if doc.is_empty:
        raise ParsingError(f"'{f.name}' produced no extractable text.")
    return doc


# --- per-type parsers ---

def _parse_pdf(f: AcceptedFile) -> ParsedDocument:
    if config.enable_docling:
        try:
            return _parse_pdf_docling(f)
        except Exception as e:
            logger.warning("Docling failed on '%s' (%s); falling back to PyMuPDF.", f.name, e)
    return _parse_pdf_pymupdf(f)


def _parse_pdf_pymupdf(f: AcceptedFile) -> ParsedDocument:
    import fitz  # PyMuPDF

    sections: List[Section] = []
    ocr_pages = 0

    with fitz.open(stream=f.data, filetype="pdf") as pdf:
        for i, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            if not text:
                # No embedded text — page is scanned/handwritten; fall back to OCR.
                text = _ocr_pdf_page(page, f.name, i)
                if text:
                    ocr_pages += 1
            if text:
                sections.append(Section(text=text, heading=f"Page {i}"))

    if ocr_pages:
        logger.info("'%s': OCR applied to %d page(s).", f.name, ocr_pages)

    return ParsedDocument(filename=f.name, sections=sections)


def _ocr_pdf_page(page, filename: str, page_num: int) -> str:
    """Render a PDF page at 300 DPI and run Tesseract OCR with preprocessing."""
    try:
        import fitz
        import pytesseract
        from PIL import Image

        if config.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = config.tesseract_cmd

        pix = page.get_pixmap(dpi=300, colorspace=fitz.csRGB)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img = _preprocess_for_ocr(img)
        return pytesseract.image_to_string(img, config="--oem 3 --psm 6").strip()
    except ImportError:
        logger.warning(
            "pytesseract not available; cannot OCR page %d of '%s'.", page_num, filename
        )
        return ""
    except Exception as e:
        logger.warning("OCR failed for page %d of '%s': %s", page_num, filename, e)
        return ""


def _preprocess_for_ocr(img):
    """Grayscale + contrast boost + binarize for Tesseract accuracy."""
    from PIL import ImageEnhance, ImageFilter

    img = img.convert("L")
    img = ImageEnhance.Contrast(img).enhance(2.0)
    img = img.filter(ImageFilter.SHARPEN)
    img = img.point(lambda x: 0 if x < 150 else 255)
    return img


def _img_to_jpeg_bytes(img) -> bytes:
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=88)
    return buf.getvalue()


def _parse_pdf_docling(f: AcceptedFile) -> ParsedDocument:
    """Opt-in Docling path for complex/scanned PDFs (heavy, lazy import)."""
    from docling.document_converter import DocumentConverter  # type: ignore

    src = io.BytesIO(f.data)
    converter = DocumentConverter()
    result = converter.convert(src, filename=f.name)
    markdown = result.document.export_to_markdown()
    return ParsedDocument(filename=f.name, sections=[Section(text=markdown)])


def _parse_docx(f: AcceptedFile) -> ParsedDocument:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(f.data))
    sections: List[Section] = []
    current_heading: Optional[str] = None
    buffer: List[str] = []

    def flush() -> None:
        if buffer:
            sections.append(Section(text="\n".join(buffer).strip(), heading=current_heading))
            buffer.clear()

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            flush()
            current_heading = text
        else:
            buffer.append(text)
    flush()

    # Tables -> tab-separated rows appended as their own section.
    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            sections.append(Section(text="\n".join(rows), heading="Table"))

    return ParsedDocument(filename=f.name, sections=sections)


def _parse_pptx(f: AcceptedFile) -> ParsedDocument:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(f.data))
    sections: List[Section] = []
    for i, slide in enumerate(prs.slides, start=1):
        title: Optional[str] = None
        lines: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title is None and shape == slide.shapes.title:
                title = text
            else:
                lines.append(text)
        heading = f"Slide {i}" + (f": {title}" if title else "")
        body = ((title + "\n") if title else "") + "\n".join(lines)
        if body.strip():
            sections.append(Section(text=body.strip(), heading=heading))
    return ParsedDocument(filename=f.name, sections=sections)


def _parse_xlsx(f: AcceptedFile) -> ParsedDocument:
    import pandas as pd

    sections: List[Section] = []
    sheets = pd.read_excel(io.BytesIO(f.data), sheet_name=None, engine="openpyxl")
    for name, df in sheets.items():
        if df.empty:
            continue
        sections.append(Section(text=df.to_csv(index=False).strip(), heading=f"Sheet: {name}"))
    return ParsedDocument(filename=f.name, sections=sections)


def _parse_csv(f: AcceptedFile) -> ParsedDocument:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(f.data))
    text = df.to_csv(index=False).strip()
    return ParsedDocument(filename=f.name, sections=[Section(text=text, heading=f.name)])


def _parse_json(f: AcceptedFile) -> ParsedDocument:
    obj = json.loads(f.data.decode("utf-8", errors="replace"))
    pretty = json.dumps(obj, indent=2, ensure_ascii=False)
    return ParsedDocument(filename=f.name, sections=[Section(text=pretty, heading=f.name)])


def _parse_text(f: AcceptedFile) -> ParsedDocument:
    try:
        text = f.data.decode("utf-8")
    except UnicodeDecodeError:
        text = f.data.decode("latin-1")
    return ParsedDocument(filename=f.name, sections=[Section(text=text.strip())])


def _parse_html(f: AcceptedFile) -> ParsedDocument:
    from bs4 import BeautifulSoup

    raw = f.data.decode("utf-8", errors="replace")
    soup = BeautifulSoup(raw, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = soup.title.string.strip() if soup.title and soup.title.string else None
    text = soup.get_text(separator="\n")
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return ParsedDocument(
        filename=f.name,
        sections=[Section(text="\n".join(lines), heading=title)],
    )


def _parse_image(f: AcceptedFile) -> ParsedDocument:
    import pytesseract
    from PIL import Image

    if config.tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = config.tesseract_cmd

    img = Image.open(io.BytesIO(f.data))
    img = _preprocess_for_ocr(img)
    text = pytesseract.image_to_string(img, config="--oem 3 --psm 6").strip()
    if not text:
        raise ParsingError(f"OCR found no text in image '{f.name}'.")
    return ParsedDocument(filename=f.name, sections=[Section(text=text)])
